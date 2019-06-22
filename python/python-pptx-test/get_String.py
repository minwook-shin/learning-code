from pptx import Presentation

presentation = Presentation("test2.pptx")

text_runs = []

for slide in presentation.slides:
    for shape in slide.shapes:
        for run in shape.text_frame.paragraphs:
            text_runs.append(run.text)


print(text_runs)