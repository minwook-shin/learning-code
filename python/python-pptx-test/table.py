from pptx import Presentation
from pptx.util import Inches

pesentation = Presentation()
title_only_slide_layout = pesentation.slide_layouts[5]
slide = pesentation.slides.add_slide(title_only_slide_layout)

shapes = slide.shapes
shapes.title.text = 'Adding Table'

table = shapes.add_table(rows=2, cols=2, left=Inches(2.0), top=Inches(2.0), width=Inches(6.0), height=Inches(1.0)).table

table.cell(0, 0).text = 'name'
table.cell(0, 1).text = 'desc'

table.cell(1, 0).text = 'test_name'
table.cell(1, 1).text = 'test desc'

pesentation.save('test4.pptx')
