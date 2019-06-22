from pptx import Presentation

presentation = Presentation()
bullet_slide_layout = presentation.slide_layouts[1]
slide = presentation.slides.add_slide(bullet_slide_layout)

shapes = slide.shapes
title_shape = shapes.title
title_shape.text = 'Adding Slide'

body_shape = shapes.placeholders[1]
tf = body_shape.text_frame
tf.text = 'bullet slide layout'

p = tf.add_paragraph()
p.text = 'add 1 level paragraph'
p.level = 1

p = tf.add_paragraph()
p.text = 'add 2 level paragraph'
p.level = 2

presentation.save('test2.pptx')