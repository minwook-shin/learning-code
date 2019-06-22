from pptx import Presentation
from pptx.util import Inches

img_path = 'test.png'

presentation = Presentation()
blank_slide_layout = presentation.slide_layouts[6]
slide = presentation.slides.add_slide(blank_slide_layout)

slide.shapes.add_picture(img_path, Inches(1), Inches(1))

presentation.save('test3.pptx')